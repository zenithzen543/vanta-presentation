"""Generate the Thorn & Furrow portfolio Word doc and PowerPoint (10 min + NDA appendix)."""

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Inches, Pt, RGBColor
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

OUT = Path(__file__).resolve().parent
INK = RGBColor(0x12, 0x11, 0x0F)
MOSS = RGBColor(0x3A, 0x4A, 0x38)
MUTED = RGBColor(0x6B, 0x65, 0x60)
INK_PPT = PptRGB(0x12, 0x11, 0x0F)
MOSS_PPT = PptRGB(0x3A, 0x4A, 0x38)
PAPER_PPT = PptRGB(0xF6, 0xF4, 0xF0)
WHITE_PPT = PptRGB(0xFF, 0xFF, 0xFF)
MUTED_PPT = PptRGB(0x6B, 0x65, 0x60)
RULE_PPT = PptRGB(0xE4, 0xDF, 0xD6)
SAGE_PPT = PptRGB(0xA8, 0xB5, 0xA4)
CREAM_PPT = PptRGB(0xD8, 0xD4, 0xCC)

BEATS = [
    {
        "clock": "0:00–0:25",
        "slide": "1 · Title",
        "do": "Deck on the projector. Catalog open at localhost:3000. Studio logged in on a second screen. Do not demo yet.",
        "say": (
            "I am presenting a marketing website — not a SaaS product. "
            "Next.js is the site. Sanity is the CMS. Vercel is where it deploys. "
            "This catalog is Thorn and Furrow. I can open it. "
            "Two other live sites come after the close. I will not open their studios."
        ),
    },
    {
        "clock": "0:25–1:05",
        "slide": "2 · Not a SaaS",
        "do": "Point at the two columns. SaaS on the left. This site on the right. Stay on the slide.",
        "say": (
            "A normal SaaS is an app people log into. Dashboard, settings, roles, billing. The software is the product. "
            "This is not that. Nobody logs in. There is no account. "
            "This is a public marketing site for a real business — a farm that sells seed. "
            "The pages are the catalog: packets, stories, a letter to order. "
            "Marketers own the words and the stock. Engineering owns the types so those changes do not need a deploy. "
            "The conversion is a packet-list request — a lead — not a SaaS signup."
        ),
    },
    {
        "clock": "1:05–1:50",
        "slide": "3 · The stack",
        "do": "Walk the three boxes: Next.js, then Sanity, then Vercel. Do not open code.",
        "say": (
            "Three tools. Next.js is the website — React, App Router, what visitors see. "
            "Sanity is the CMS. Studio is the editor. Packets and homepage sections are documents, not a Word file. "
            "Vercel is the host. Git push, a live URL, CDN, ISR. That is production for a Next.js site. "
            "Editor publishes in Sanity. Next.js pulls the page. Vercel serves it. "
            "That loop is what I am here to show."
        ),
    },
    {
        "clock": "1:50–2:25",
        "slide": "4 · The job",
        "do": "Stay on the pitch slide. Do not click the site yet.",
        "say": (
            "The farm packs seed March through June. This season is Catalog Number 14. "
            "If a variety is listed, they grow it. "
            "Marketing does not wait on engineering for a new heading or a sold-out packet. Those are fields. "
            "I own the types. They own the weekly catalog. "
            "That is the same contract as any marketing site: structured content, page composition, caching, SEO, "
            "and analytics that survive a copy change."
        ),
    },
    {
        "clock": "2:25–3:00",
        "slide": "5 · What visitors see",
        "do": "Point at the four cards, then glance at the live homepage if it is visible.",
        "say": (
            "Visitors land on a homepage composed from Sanity sections — masthead, featured packet, sowing table, the spring list. "
            "Varieties is the catalog: seventeen packets grouped by family, the way a print catalog is grouped. "
            "Journal is the seasonal notebook. Mail order is a letter, not a cart: mark packets, tell us your zone, we write back. "
            "Studio exists at /studio for editors. It is not linked on the public site. The storefront should feel live, not like a CMS demo."
        ),
    },
    {
        "clock": "3:00–3:35",
        "slide": "6 · Why this shape",
        "do": None,
        "say": (
            "I did not build a hero, a two-by-three feature grid, and a logo wall. That would look like homework. "
            "A catalog needs four things a marketing site actually needs. "
            "Documents: Icehouse Tomato is one record, reused on home, the index, the packet page, and the order form. "
            "Composition: the homepage is an ordered list of blocks. Editors reorder it without a deploy. "
            "Editor-safe fields: stock, packet price, sowing, isolation, photographs. "
            "Stable tracking: GTM keys off event IDs, not the words on the button."
        ),
    },
    {
        "clock": "3:35–4:25",
        "slide": "7 · Architecture",
        "do": "Walk the five boxes left to right. Do not open code unless asked.",
        "say": (
            "The path is simple. An editor publishes in Sanity Studio. "
            "Next.js pulls the page with GROQ through next-sanity. "
            "Pages are React Server Components in the App Router. "
            "Production uses ISR — pages revalidate about every sixty seconds — plus a signed revalidate API after publish, so we do not wait a full minute. "
            "If Sanity is down, a local seed still renders the catalog. That is resilience, and it is also how a demo does not die. "
            "Photographs come from Sanity assets, with a local file fallback per slug."
        ),
    },
    {
        "clock": "4:25–5:00",
        "slide": "8 · Content model",
        "do": "Name the five documents. Mention the Studio desk is labeled The farm.",
        "say": (
            "Five document types. Farm settings: name, nav, primary CTA, footer, default SEO. "
            "Catalog homepage: SEO plus an ordered section array. "
            "Variety: latin name, family, days, packet, stock, photo, story. That is the product object. "
            "Journal note and workshop are seasonal content. "
            "The Studio desk is labeled The farm so editors see a catalog, not a generic CMS. "
            "Change Icehouse Tomato once. Home, packet page, and order form all update, because they reference the same document."
        ),
    },
    {
        "clock": "5:00–7:40",
        "slide": "9 · Live demo",
        "do": (
            "Switch to the browser. (1) Homepage: masthead, Icehouse feature, sowing table, packet grid. "
            "(2) Click Icehouse Tomato. (3) Studio → Catalog homepage → masthead heading → Publish → reload /. "
            "(4) Hover Request a packet list and say the event ID."
        ),
        "say": (
            "This is the site. Catalog Number 14, Spring 2026. The masthead is a section. "
            "The featured packet is a document — latin name, story, maturity, stock — not a dashboard widget. "
            "The sowing table is editorial and dated for zone 6a. "
            "The spring list is seventeen references to variety documents. Marketing reorders this list. I own the types. "
            "Icehouse Tomato: same document you just saw featured. App Router, static params, ISR. "
            "In Studio, under Catalog homepage, I change the masthead heading and publish. Reload. The heading is live. That publish loop is the work. "
            "Request a packet list: the label is copy. The contract is eventId — hero_order, nav_order. Copy can change. GTM does not."
        ),
    },
    {
        "clock": "7:40–8:20",
        "slide": "10 · What I would defend",
        "do": "Back to the deck. If Studio failed, say: seed fallback still rendered the catalog.",
        "say": (
            "Five decisions I would make again. "
            "Server components and GROQ, not a client-side CMS fetch on every page. "
            "References over duplication. "
            "Event IDs on every CTA so measurement survives a rewrite. "
            "ISR plus a signed revalidate webhook. "
            "No cart. The conversion is a qualified packet-list request — closer to a real marketing funnel than a fake checkout."
        ),
    },
    {
        "clock": "8:20–8:50",
        "slide": "11 · Close",
        "do": "One line. If they only have ten minutes, stop and take questions. Otherwise go to slide 12.",
        "say": (
            "The catalog is the product. The CMS is how it changes without a deploy."
        ),
    },
    {
        "clock": "8:50–9:05",
        "slide": "12 · NDA map",
        "do": "Point at both public homepages. Do not open Studio. Then Jamb, then Slingshot.",
        "say": (
            "Same pattern as an ads business: who buys, what is inventory, how a visit becomes a sale. Then where the CMS sits. "
            "Jamb is a Pimlico dealer. Slingshot Bio sells reagents to flow cytometry labs. "
            "I will not open Studio, schemas, or fields. Public surfaces only."
        ),
    },
    {
        "clock": "9:05–9:40",
        "slide": "13 · Jamb business",
        "do": "Walk who, inventory, sale. Screenshot is the public homepage. No Studio.",
        "say": (
            "Jamb opened on Pimlico Road in 2001. The buyer is an interior designer, a collector, a country-house client — London, then Los Angeles, Chicago, Dallas, Atlanta, Palm Beach. "
            "Inventory is two things at once. Unique antiques: chimneypieces from the seventeenth century on, furniture, lighting. "
            "And handmade reproductions from the south London workshop, so a design stays available after the antique sells — that is the original business idea. "
            "The sale is high-ticket and consultative. A designer specs a piece from the catalog. The showroom closes it. Enquiry, not Amazon checkout. "
            "Sanity owns collections and journal stories. Next.js is the storefront."
        ),
    },
    {
        "clock": "9:40–10:10",
        "slide": "14 · Jamb catalog",
        "do": "Name the four public collections. Do not invent schemas.",
        "say": (
            "Four public surfaces, same contract as any marketing catalog. "
            "Fireplaces: antique chimneypieces and stone or marble reproductions, plus a bespoke service. "
            "Lighting: hanging globes, lanterns, wall lights — the Original Globe is the story that started the reproduction line. "
            "Furniture: English country-house seating and tables, antique and made. "
            "Journal: stories for architects and designers, not a blog calendar. "
            "Editors change stock and stories. Engineering owns the types."
        ),
    },
    {
        "clock": "10:10–10:45",
        "slide": "15 · Slingshot business",
        "do": "Walk the problem, then the product, then who pays. Screenshot is How Mimics Work. No Studio.",
        "say": (
            "Slingshot Bio sells to flow cytometry labs: biopharma, CROs, academic cores, cell therapy, instrument makers. "
            "The problem is the control. Donor cells expire and drift lot to lot. Polystyrene beads have the wrong scatter. "
            "They sell shelf-stable synthetic cell mimics — polymer particles engineered to scatter, fluoresce, and stain like real cells. "
            "Compensation, unmixing, viability, immunophenotyping, instrument setup. Catalog SKUs plus custom mimics on request. "
            "A scientist finds a control, reads a protocol, orders a reagent. Repeat purchases. "
            "Sanity owns product and resource documents. Next.js is the catalog and shop."
        ),
    },
    {
        "clock": "10:45–11:15",
        "slide": "16 · Slingshot catalog",
        "do": "Home, how it works, shop. Name public surfaces only.",
        "say": (
            "Three public surfaces. Homepage is the catalog of controls. "
            "How Mimics Work is the argument: not a biologic, not a bead. "
            "Shop is the conversion: SKUs grouped by purpose — unmixing, immunophenotyping, CAR-T, custom — then a cart. "
            "Resources sit beside the shop: application notes, data sheets, protocols. That is the education that sells a reagent. "
            "Same rule: I cannot open the desk."
        ),
    },
    {
        "clock": "11:15–11:30",
        "slide": "17 · Stop",
        "do": "Stop. Invite questions on Thorn and Furrow, not on those studios.",
        "say": (
            "I cannot open those desks. This catalog I can. "
            "Happy to go into schema, caching, or the order form — here."
        ),
    },
]


def set_run(run, *, size=11, bold=False, color=INK, italic=False, name="Calibri"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn("w:eastAsia"), name)


def add_para(doc, text, *, size=11, bold=False, italic=False, color=INK, name="Calibri", space_before=0, space_after=8, align=None):
    p = doc.add_paragraph()
    if align:
        p.alignment = align
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.line_spacing = 1.15
    run = p.add_run(text)
    set_run(run, size=size, bold=bold, italic=italic, color=color, name=name)
    return p


def add_heading(doc, text, level=1):
    return add_para(
        doc,
        text,
        size=18 if level == 1 else 13,
        bold=True,
        color=MOSS if level == 1 else INK,
        name="Georgia",
        space_before=18 if level == 1 else 12,
        space_after=8,
    )


def shade_cell(cell, hex_color):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tcPr.append(shd)


def set_cell_text(cell, text, *, bold=False, color=INK, size=10):
    cell.text = ""
    run = cell.paragraphs[0].add_run(text)
    set_run(run, size=size, bold=bold, color=color)


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        set_cell_text(cell, header, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), size=10)
        shade_cell(cell, "3A4A38")
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.rows[r].cells[c]
            set_cell_text(cell, value, size=10)
            if r % 2 == 0:
                shade_cell(cell, "F6F4F0")
    doc.add_paragraph()


def build_docx():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.85)
    section.bottom_margin = Inches(0.85)
    section.left_margin = Inches(1.0)
    section.right_margin = Inches(1.0)

    add_para(doc, "10-MINUTE PRESENTATION  ·  SPEAKING SCRIPT", size=10, bold=True, color=MOSS, space_after=4)
    add_para(doc, "Thorn & Furrow", size=28, bold=True, name="Georgia", space_after=2)
    add_para(
        doc,
        "Read the SAY lines out loud. Follow DO. Stay on the clock. Ten minutes is slides 1–11. Slides 12–17 if they have time.",
        size=12,
        italic=True,
        color=MUTED,
        name="Georgia",
        space_after=10,
    )

    add_heading(doc, "Before you start")
    for item in [
        "Two windows: catalog at http://localhost:3000 and Studio at http://localhost:3000/studio (already logged in).",
        "Hard-refresh the catalog once. Confirm packets and photos load.",
        "Do not mention NDA, interview scaffolding, or the CMS on the public footer — the site is a live catalog.",
        "If Studio fails: keep going on the public site. Say the seed fallback still renders the catalog.",
        "If you run long, skip slide 10 (decisions), close, and take questions. Slides 12–17 are extra if they ask about past work.",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        set_run(run, size=11)

    add_heading(doc, "Clock")
    add_table(
        doc,
        ["Time", "Slide", "You are doing"],
        [
            ["0:00–0:25", "1 Title", "Marketing site, not SaaS. This catalog I can open."],
            ["0:25–1:05", "2 Not a SaaS", "Login app vs public catalog. Lead, not signup."],
            ["1:05–1:50", "3 The stack", "Next.js = site. Sanity = CMS. Vercel = deploy."],
            ["1:50–2:25", "4 The job", "Editors vs engineering. Same contract as a marketing site."],
            ["2:25–3:00", "5 Product", "Four surfaces. No cart."],
            ["3:00–3:35", "6 Why this shape", "Documents, composition, fields, event IDs."],
            ["3:35–4:25", "7 Architecture", "Studio → GROQ → Next.js → ISR."],
            ["4:25–5:00", "8 Content model", "Five documents. One variety, many pages."],
            ["5:00–7:40", "9 Live demo", "Home, packet, publish, tracking. This is the talk."],
            ["7:40–8:20", "10 Decisions", "Five calls you would make again."],
            ["8:20–8:50", "11 Close", "One line. Stop here if they only have ten minutes."],
            ["8:50–9:05", "12 NDA map", "Jamb and Slingshot. Public only. No Studio."],
            ["9:05–9:40", "13 Jamb business", "Who buys, inventory, how a visit becomes a sale."],
            ["9:40–10:10", "14 Jamb catalog", "Fireplaces, lighting, furniture, journal."],
            ["10:10–10:45", "15 Slingshot business", "Problem, product, who pays, conversion."],
            ["10:45–11:15", "16 Slingshot catalog", "Home, how it works, shop, resources."],
            ["11:15–11:30", "17 Stop", "I cannot open those desks. Questions here."],
        ],
    )

    add_heading(doc, "Speaking script")
    add_para(
        doc,
        "Read SAY. Do not improvise new features. If they interrupt, finish the sentence, then answer.",
        italic=True,
        color=MUTED,
        space_after=12,
    )

    for beat in BEATS:
        add_para(doc, f"{beat['clock']}    {beat['slide']}", size=13, bold=True, color=MOSS, name="Georgia", space_before=14, space_after=4)
        if beat.get("do"):
            add_para(doc, f"DO  ·  {beat['do']}", size=10, italic=True, color=MUTED, space_after=6)
        add_para(doc, "SAY", size=9, bold=True, color=MOSS, space_after=2)
        add_para(doc, beat["say"], size=12, space_after=10)

    add_heading(doc, "If they ask")
    add_table(
        doc,
        ["Question", "Answer in one breath"],
        [
            ["Is this a SaaS?", "No. Nobody logs in. It is a public marketing catalog. The conversion is a lead, not a subscription."],
            ["Can you show Jamb or Slingshot Studio?", "No. NDA. I can name the business and which public surfaces Sanity feeds. Thorn & Furrow is the walkthrough."],
            ["How does the homepage work?", "An ordered Sanity array. The UI switches on section _type. Editors reorder blocks."],
            ["What is ISR?", "Pages regenerate on a 60-second window, plus a signed revalidate call after publish."],
            ["How do you track CTAs?", "data-event-id on the link. GTM listens for the ID, not the label."],
            ["What if Sanity is down?", "Local seed still renders. The catalog does not go blank."],
            ["Is there a cart?", "No. Mail order is a qualified packet-list request — a marketing lead."],
        ],
    )

    add_heading(doc, "Close line (memorize)")
    add_para(
        doc,
        "The catalog is the product. The CMS is how it changes without a deploy.",
        size=14,
        italic=True,
        name="Georgia",
        space_after=18,
    )

    add_para(doc, "Leave-behind  ·  one page", size=10, bold=True, color=MOSS, space_before=8, space_after=6)
    add_para(
        doc,
        "Thorn & Furrow is the walkable catalog: 17 varieties, journal, workshops, mail-order letter. "
        "Jamb (Pimlico dealer: antiques + reproductions; showroom sale) and Slingshot Bio (cell-mimic reagents for flow cytometry labs; catalog + shop) are live Next.js + Sanity sites under NDA — business and public CMS surfaces only. "
        "Homepage is an ordered section array. Varieties are documents referenced on home, index, packet page, and order form. "
        "Stack: TypeScript, Tailwind, ISR, GROQ, event IDs (nav_order, hero_order, footer_order, variety_order). Studio desk: The farm.",
        size=11,
        space_after=8,
    )
    add_para(
        doc,
        "Stack  ·  Next.js (App Router)  ·  Sanity  ·  TypeScript  ·  Tailwind CSS  ·  Vercel-ready",
        size=9,
        color=MUTED,
        space_before=8,
    )

    path = OUT / "Thorn-and-Furrow-Portfolio.docx"
    doc.save(path)
    return path


SANS = "Calibri"
SERIF = "Georgia"
TOTAL = 17
MARGIN = PptInches(0.72)
SLIDE_W = PptInches(13.333)
SLIDE_H = PptInches(7.5)
CREAM_BG = PptRGB(0xED, 0xE8, 0xDC)
GHOST_PPT = PptRGB(0x2A, 0x36, 0x28)
TOMATO = Path(__file__).resolve().parents[1] / "public" / "varieties" / "icehouse-tomato.jpg"


def set_fill(shape, color, alpha=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is None:
        return shape
    solid = shape._element.spPr.find(f"{{{A_NS}}}solidFill")
    srgb = solid.find(f"{{{A_NS}}}srgbClr") if solid is not None else None
    if srgb is None:
        return shape
    for child in list(srgb):
        if child.tag == f"{{{A_NS}}}alpha":
            srgb.remove(child)
    etree.SubElement(srgb, f"{{{A_NS}}}alpha").set("val", str(int(alpha * 100000)))
    return shape


def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    set_fill(shape, color, alpha=alpha)
    return shape


def add_round(slide, l, t, w, h, fill, line=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = PptPt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    set_fill(shape, fill)
    return shape


def hairline(slide, l, t, w, color=RULE_PPT):
    return add_rect(slide, l, t, w, PptPt(1.15), color)


def make_cover(src: Path, dest: Path, w: int, h: int) -> Path:
    im = Image.open(src).convert("RGB")
    sw, sh = im.size
    target = w / h
    src_a = sw / sh
    if src_a > target:
        nw = int(sh * target)
        x = (sw - nw) // 2
        im = im.crop((x, 0, x + nw, sh))
    else:
        nh = int(sw / target)
        y = (sh - nh) // 2
        im = im.crop((0, y, sw, y + nh))
    im = im.resize((w, h), Image.Resampling.LANCZOS)
    dest.parent.mkdir(parents=True, exist_ok=True)
    im.save(dest, "JPEG", quality=92, optimize=True)
    return dest


SHOTS = OUT / "shots"
FITTED = OUT / "_fitted"
CHROME_BG = PptRGB(0xEE, 0xEA, 0xE4)
DOT_RED = PptRGB(0xC4, 0x7C, 0x6E)
DOT_GOLD = PptRGB(0xD2, 0xBE, 0x7A)
DOT_GREEN = PptRGB(0x8A, 0xA0, 0x82)


def emu_in(x):
    return int(x) / 914400.0


def fit_top(src: Path, dest: Path, w: int, h: int) -> Path:
    im = Image.open(src).convert("RGB")
    sw, sh = im.size
    target = w / h
    src_a = sw / sh
    if src_a > target:
        nw = int(sh * target)
        x = (sw - nw) // 2
        im = im.crop((x, 0, x + nw, sh))
    else:
        nh = min(sh, int(sw / target))
        im = im.crop((0, 0, sw, nh))
    im = im.resize((w, h), Image.Resampling.LANCZOS)
    dest.parent.mkdir(parents=True, exist_ok=True)
    im.save(dest, "JPEG", quality=86, optimize=True)
    return dest


def add_browser_shot(slide, src, l, t, w, h, url=""):
    add_round(slide, l, t, w, h, WHITE_PPT, RULE_PPT, 0.04)
    chrome_h = PptInches(0.3) if url else PptInches(0)
    if url:
        add_rect(slide, l, t, w, chrome_h, CHROME_BG)
        for i, color in enumerate((DOT_RED, DOT_GOLD, DOT_GREEN)):
            add_oval(slide, l + PptInches(0.14 + i * 0.18), t + PptInches(0.09), PptInches(0.12), PptInches(0.12), color)
        add_textbox(
            slide,
            l + PptInches(0.72),
            t + PptInches(0.04),
            w - PptInches(0.9),
            PptInches(0.24),
            url,
            size=10,
            color=MUTED_PPT,
        )
    inner_t = t + chrome_h
    inner_h = h - chrome_h
    px_w = max(480, int(emu_in(w) * 140))
    px_h = max(280, int(emu_in(inner_h) * 140))
    fitted = fit_top(Path(src), FITTED / f"{Path(src).stem}_{px_w}x{px_h}.jpg", px_w, px_h)
    slide.shapes.add_picture(str(fitted), l, inner_t, w, inner_h)


def style_run(run, *, size, bold=False, color=INK_PPT, font=SANS):
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, l, t, w, h, text, *, size=18, bold=False, color=INK_PPT, font=SANS, align=PP_ALIGN.LEFT, after=0):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    parts = str(text).split("\n")
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = PptPt(after)
        run = p.add_run()
        run.text = part
        style_run(run, size=size, bold=bold, color=color, font=font)
    return box


def num_mark(slide, l, t, num):
    add_round(slide, l, t, PptInches(0.44), PptInches(0.44), MOSS_PPT, None, 0.18)
    add_textbox(
        slide,
        l,
        t + PptInches(0.06),
        PptInches(0.44),
        PptInches(0.34),
        num,
        size=11,
        bold=True,
        color=WHITE_PPT,
        align=PP_ALIGN.CENTER,
    )


def add_lines(slide, l, t, w, h, lines, *, size=16, color=INK_PPT, font=SANS, after=8, align=PP_ALIGN.LEFT, bold=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = PptPt(after)
        run = p.add_run()
        run.text = line
        style_run(run, size=size, bold=bold, color=color, font=font)
    return box


def notes(slide, beat):
    body = f"{beat['clock']}  ·  {beat['slide']}\n\n"
    if beat.get("do"):
        body += f"DO: {beat['do']}\n\n"
    body += f"SAY:\n{beat['say']}"
    slide.notes_slide.notes_text_frame.text = body


def footer(slide, index, dark=False):
    muted = SAGE_PPT if dark else MUTED_PPT
    rule = PptRGB(0x2A, 0x36, 0x28) if dark else RULE_PPT
    hairline(slide, MARGIN, PptInches(6.98), PptInches(11.9), rule)
    add_textbox(slide, MARGIN, PptInches(7.08), PptInches(4.2), PptInches(0.28), "Thorn & Furrow", size=10, color=muted, font=SERIF)
    add_textbox(slide, PptInches(9.0), PptInches(7.08), PptInches(3.6), PptInches(0.28), f"{index:02d}  /  {TOTAL:02d}", size=10, color=muted, align=PP_ALIGN.RIGHT)


def blank(prs, bg=PAPER_PPT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
    return slide


def content_slide(prs, title, index):
    slide = blank(prs)
    add_rect(slide, 0, 0, PptInches(0.1), SLIDE_H, MOSS_PPT)
    add_rect(slide, 0, 0, SLIDE_W, PptInches(0.07), MOSS_PPT)
    add_textbox(slide, MARGIN, PptInches(0.32), PptInches(12), PptInches(0.72), title, size=32, bold=True, color=INK_PPT, font=SERIF)
    hairline(slide, MARGIN, PptInches(1.12), PptInches(11.9), RULE_PPT)
    footer(slide, index)
    return slide


def build_pptx(path=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Thorn & Furrow"
    prs.core_properties.subject = "A live heirloom-seed catalog on Next.js and Sanity"

    # 1 Title — magazine spread: ink + Icehouse Tomato
    panel_l = PptInches(8.02)
    panel_w = SLIDE_W - panel_l
    cover = make_cover(TOMATO, OUT / "_title-panel.jpg", 1200, 1620)
    s = blank(prs, INK_PPT)
    s.shapes.add_picture(str(cover), panel_l, 0, panel_w, SLIDE_H)
    add_rect(s, 0, 0, panel_l, SLIDE_H, INK_PPT)
    add_rect(s, panel_l, 0, PptInches(0.09), SLIDE_H, MOSS_PPT)
    add_rect(s, 0, 0, panel_l, PptInches(0.07), MOSS_PPT)
    add_textbox(s, PptInches(4.6), PptInches(0.35), PptInches(3.2), PptInches(1.5), "14", size=96, bold=True, color=GHOST_PPT, font=SERIF, align=PP_ALIGN.RIGHT)
    add_round(s, MARGIN, PptInches(1.42), PptInches(4.15), PptInches(0.38), MOSS_PPT, None, 0.2)
    add_textbox(s, MARGIN, PptInches(1.48), PptInches(4.15), PptInches(0.28), "TEN MINUTES  ·  THEN PAST WORK", size=10, bold=True, color=WHITE_PPT, align=PP_ALIGN.CENTER)
    add_textbox(s, MARGIN, PptInches(2.05), PptInches(7.1), PptInches(1.2), "Thorn & Furrow", size=44, bold=True, color=WHITE_PPT, font=SERIF)
    add_textbox(
        s,
        MARGIN,
        PptInches(3.3),
        PptInches(6.9),
        PptInches(1.6),
        "A live heirloom-seed catalog\non Next.js and Sanity.",
        size=22,
        color=CREAM_PPT,
        font=SERIF,
        after=6,
    )
    add_textbox(s, MARGIN, PptInches(5.15), PptInches(6.9), PptInches(0.4), "Product.  Content model.  Publish loop.", size=16, color=SAGE_PPT)
    add_rect(s, MARGIN, PptInches(6.35), PptInches(1.45), PptInches(0.07), MOSS_PPT)
    add_textbox(s, MARGIN, PptInches(6.55), PptInches(6.9), PptInches(0.55), "Next.js   ·   Sanity   ·   Vercel", size=13, color=WHITE_PPT)
    notes(s, BEATS[0])

    # 2 Not a SaaS
    s = content_slide(prs, "This is not a SaaS", 2)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.28),
        PptInches(11.9),
        PptInches(0.48),
        "I am presenting a public marketing website for a real business. Visitors do not log in. Marketers change the catalog without a deploy.",
        size=15,
        color=MUTED_PPT,
    )
    add_round(s, MARGIN, PptInches(1.88), PptInches(5.85), PptInches(4.78), WHITE_PPT, RULE_PPT, 0.04)
    add_rect(s, MARGIN, PptInches(1.88), PptInches(5.85), PptInches(0.52), MOSS_PPT)
    add_textbox(s, PptInches(0.95), PptInches(1.98), PptInches(5.4), PptInches(0.36), "A normal SaaS", size=16, bold=True, color=WHITE_PPT, font=SERIF)
    saas_lines = [
        "People log in. Dashboard, settings, roles, billing.",
        "The software is the product.",
        "Data belongs to each customer’s account.",
        "Engineering ships features: auth, CRUD, permissions.",
        "Success is signup, activation, subscription.",
    ]
    for i, line in enumerate(saas_lines):
        add_textbox(s, PptInches(0.95), PptInches(2.55) + i * PptInches(0.72), PptInches(5.35), PptInches(0.68), line, size=14, color=MUTED_PPT)
    add_round(s, PptInches(6.78), PptInches(1.88), PptInches(5.85), PptInches(4.78), WHITE_PPT, RULE_PPT, 0.04)
    add_rect(s, PptInches(6.78), PptInches(1.88), PptInches(5.85), PptInches(0.52), MOSS_PPT)
    add_textbox(s, PptInches(7.0), PptInches(1.98), PptInches(5.4), PptInches(0.36), "This site", size=16, bold=True, color=WHITE_PPT, font=SERIF)
    site_lines = [
        "Nobody logs in. The catalog is public.",
        "The catalog is the marketing surface.",
        "Content is editorial: packets, stories, stock.",
        "Engineering ships types, pages, publish, SEO, tracking.",
        "Success is a packet-list request — a lead.",
    ]
    for i, line in enumerate(site_lines):
        add_textbox(s, PptInches(7.0), PptInches(2.55) + i * PptInches(0.72), PptInches(5.35), PptInches(0.68), line, size=14, color=MUTED_PPT)
    notes(s, BEATS[1])

    # 3 The stack
    s = content_slide(prs, "Next.js, Sanity, Vercel", 3)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.28),
        PptInches(11.9),
        PptInches(0.4),
        "Three tools. Editor publishes in Sanity. Next.js draws the page. Vercel hosts it.",
        size=15,
        color=MUTED_PPT,
    )
    stack = [
        ("01", "Next.js", "The website. React, App Router. What visitors see: homepage, packets, journal, mail order."),
        ("02", "Sanity", "The CMS. Studio is the editor. Packets and homepage sections are documents — not a Word file."),
        ("03", "Vercel", "The host. Git push, a live URL, CDN, ISR. Production for a Next.js marketing site."),
    ]
    for i, (num, name, body) in enumerate(stack):
        left = MARGIN + i * PptInches(4.05)
        add_round(s, left, PptInches(1.85), PptInches(3.9), PptInches(4.8), WHITE_PPT, RULE_PPT, 0.05)
        add_oval(s, left + PptInches(0.22), PptInches(2.08), PptInches(0.44), PptInches(0.44), MOSS_PPT)
        add_textbox(s, left + PptInches(0.22), PptInches(2.16), PptInches(0.44), PptInches(0.32), num, size=12, bold=True, color=WHITE_PPT, align=PP_ALIGN.CENTER)
        add_textbox(s, left + PptInches(0.22), PptInches(2.7), PptInches(3.45), PptInches(0.5), name, size=24, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, left + PptInches(0.22), PptInches(3.35), PptInches(3.45), PptInches(2.9), body, size=15, color=MUTED_PPT)
    notes(s, BEATS[2])

    # 4 The job — copy left, homepage right
    s = content_slide(prs, "I own the types. They own the catalog.", 4)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.42),
        PptInches(5.4),
        PptInches(1.85),
        "The farm packs seed March through June. Catalog No. 14 — if a variety is listed, they grow it.\n\nMarketing does not wait on engineering for a new heading or a sold-out packet. Those are fields. I own the types. They own the weekly catalog.",
        size=14,
        color=MUTED_PPT,
        after=6,
    )
    facts = [
        ("No. 14", "Spring 2026. The season is a document, not a deploy."),
        ("17", "Packets still grown out. No ornamental fillers."),
        ("Mar–Jun", "Packed the week they ship. Zone 6a, Tivoli."),
    ]
    for i, (stat, label) in enumerate(facts):
        top = PptInches(3.4) + i * PptInches(1.08)
        add_round(s, MARGIN, top, PptInches(5.4), PptInches(0.96), WHITE_PPT, RULE_PPT, 0.05)
        add_rect(s, MARGIN, top, PptInches(0.1), PptInches(0.96), MOSS_PPT)
        add_textbox(s, PptInches(1.05), top + PptInches(0.1), PptInches(4.75), PptInches(0.36), stat, size=20, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, PptInches(1.05), top + PptInches(0.48), PptInches(4.75), PptInches(0.38), label, size=12, color=MUTED_PPT)
    add_browser_shot(s, SHOTS / "home.png", PptInches(6.28), PptInches(1.42), PptInches(6.35), PptInches(5.26), "localhost:3000  ·  live catalog")
    notes(s, BEATS[3])

    # 5 Product — four live surfaces
    s = content_slide(prs, "What visitors see", 5)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.38),
        PptInches(11.9),
        PptInches(0.38),
        "Four public surfaces. Studio lives at /studio and is not linked on the storefront — the catalog should feel live, not like a CMS demo.",
        size=13,
        color=MUTED_PPT,
    )
    surfaces = [
        (SHOTS / "home.png", "01", "/", "Homepage", "Ordered CMS sections: masthead, featured packet, sowing table, spring list."),
        (SHOTS / "varieties.png", "02", "/varieties", "Catalog", "Seventeen packets grouped by family, the way a print catalog is grouped."),
        (SHOTS / "journal.png", "03", "/journal", "Notebook", "Seasonal field notes. Not a blog calendar."),
        (SHOTS / "order.png", "04", "/order", "Mail order", "A letter, not a cart. Mark packets, tell us your zone, we write back."),
    ]
    for i, (src, num, route, title, blurb) in enumerate(surfaces):
        col, row = i % 2, i // 2
        left = MARGIN + col * PptInches(6.1)
        top = PptInches(1.78) + row * PptInches(2.48)
        add_textbox(s, left, top, PptInches(0.7), PptInches(0.24), num, size=11, bold=True, color=MOSS_PPT)
        add_textbox(s, left + PptInches(0.7), top, PptInches(2.6), PptInches(0.24), title, size=13, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, left + PptInches(3.4), top, PptInches(2.4), PptInches(0.24), route, size=11, color=MUTED_PPT, align=PP_ALIGN.RIGHT)
        add_browser_shot(s, src, left, top + PptInches(0.26), PptInches(5.85), PptInches(1.72), "")
        add_textbox(s, left, top + PptInches(2.02), PptInches(5.85), PptInches(0.38), blurb, size=11, color=MUTED_PPT)
    notes(s, BEATS[4])

    # 6 Why — principles + packet page
    s = content_slide(prs, "Why this shape", 6)
    tiles = [
        ("01", "Documents", "Icehouse Tomato is one record. Home, the index, the packet page, and the order form all reference it. Change stock once — every surface updates."),
        ("02", "Composition", "The homepage is an ordered section array: masthead, featured packet, sowing table, spring list. Editors reorder blocks without a deploy."),
        ("03", "Editor fields", "Stock, packet price, sowing, isolation, photographs. Marketing edits fields. Engineering owns the types."),
        ("04", "Stable tracking", "GTM keys off eventId (hero_order, nav_order), not the words on the button. Copy can change. Measurement does not."),
    ]
    for i, (num, title, body) in enumerate(tiles):
        top = PptInches(1.48) + i * PptInches(1.28)
        add_round(s, MARGIN, top, PptInches(5.45), PptInches(1.16), WHITE_PPT, RULE_PPT, 0.05)
        add_textbox(s, PptInches(0.95), top + PptInches(0.08), PptInches(1.0), PptInches(0.24), num, size=11, bold=True, color=MOSS_PPT)
        add_textbox(s, PptInches(1.7), top + PptInches(0.06), PptInches(4.1), PptInches(0.28), title, size=16, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, PptInches(0.95), top + PptInches(0.38), PptInches(4.9), PptInches(0.72), body, size=12, color=MUTED_PPT)
    add_textbox(s, PptInches(6.38), PptInches(1.48), PptInches(6.2), PptInches(0.28), "Same Icehouse document — packet page", size=12, bold=True, color=MOSS_PPT)
    add_browser_shot(s, SHOTS / "icehouse.png", PptInches(6.38), PptInches(1.78), PptInches(6.22), PptInches(4.88), "/varieties/icehouse-tomato")
    notes(s, BEATS[5])

    # 7 Architecture — path + editor and storefront
    s = content_slide(prs, "Editor to production", 7)
    steps = [
        ("1", "Studio"),
        ("2", "GROQ"),
        ("3", "Next.js"),
        ("4", "ISR"),
        ("5", "Catalog"),
    ]
    for i, (num, title) in enumerate(steps):
        left = PptInches(0.72) + i * PptInches(2.5)
        add_round(s, left, PptInches(1.45), PptInches(2.18), PptInches(0.7), WHITE_PPT, RULE_PPT, 0.08)
        add_oval(s, left + PptInches(0.12), PptInches(1.56), PptInches(0.46), PptInches(0.46), MOSS_PPT)
        add_textbox(s, left + PptInches(0.12), PptInches(1.64), PptInches(0.46), PptInches(0.34), num, size=12, bold=True, color=WHITE_PPT, align=PP_ALIGN.CENTER)
        add_textbox(s, left + PptInches(0.62), PptInches(1.56), PptInches(1.45), PptInches(0.48), title, size=14, bold=True, color=INK_PPT, font=SERIF)
        if i < 4:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + PptInches(2.2),
                PptInches(1.7),
                PptInches(0.28),
                PptInches(0.16),
            )
            set_fill(arrow, SAGE_PPT)
    add_textbox(
        s,
        MARGIN,
        PptInches(2.22),
        PptInches(11.9),
        PptInches(0.42),
        "Editor publishes in Studio. Next.js pulls with GROQ. RSC pages. ISR every 60s, plus a signed revalidate API after publish. If Sanity is down, local seed still renders.",
        size=13,
        color=MUTED_PPT,
    )
    add_textbox(s, MARGIN, PptInches(2.66), PptInches(5.85), PptInches(0.26), "Studio  ·  homepage section + Icehouse document", size=12, bold=True, color=MOSS_PPT)
    add_browser_shot(s, SHOTS / "studio-home.png", MARGIN, PptInches(2.92), PptInches(5.85), PptInches(3.74), "/studio  ·  featured section")
    add_textbox(s, PptInches(6.78), PptInches(2.66), PptInches(5.85), PptInches(0.26), "Storefront  ·  the same heading, live", size=12, bold=True, color=MOSS_PPT)
    add_browser_shot(s, SHOTS / "home.png", PptInches(6.78), PptInches(2.92), PptInches(5.82), PptInches(3.74), "/")
    notes(s, BEATS[6])

    # 8 Content model — list + studio
    s = content_slide(prs, "Five documents", 8)
    rows = [
        ("01", "Farm settings", "Site name, nav, primary CTA, footer, default SEO — the chrome of the catalog."),
        ("02", "Catalog homepage", "SEO plus an ordered section array. Editors compose the front page."),
        ("03", "Variety", "The product object: latin name, family, days, packet, stock, photo, story."),
        ("04", "Journal note", "Season, excerpt, body. Field notes, not a blog calendar."),
        ("05", "Workshop", "Date, place, seats. Seasonal, same desk as the packets."),
    ]
    for i, (num, name, role) in enumerate(rows):
        top = PptInches(1.42) + i * PptInches(1.02)
        add_round(s, MARGIN, top, PptInches(5.55), PptInches(0.92), WHITE_PPT, RULE_PPT, 0.05)
        num_mark(s, PptInches(0.92), top + PptInches(0.24), num)
        add_textbox(s, PptInches(1.55), top + PptInches(0.1), PptInches(4.4), PptInches(0.32), name, size=15, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, PptInches(1.55), top + PptInches(0.44), PptInches(4.4), PptInches(0.4), role, size=11, color=MUTED_PPT)
    add_textbox(s, PptInches(6.5), PptInches(1.42), PptInches(6.1), PptInches(0.42), "Icehouse Tomato open in Studio — common name, latin, slug. Change it once; home, packet page, and order form all follow.", size=12, color=MUTED_PPT)
    add_browser_shot(s, SHOTS / "studio-variety.png", PptInches(6.5), PptInches(1.88), PptInches(6.1), PptInches(4.78), "/studio  ·  Variety document")
    notes(s, BEATS[7])

    # 9 Demo — four screens you will click
    s = content_slide(prs, "Live demo — this is the talk", 9)
    add_round(s, MARGIN, PptInches(1.42), PptInches(11.9), PptInches(0.42), MOSS_PPT, None, 0.04)
    add_textbox(s, PptInches(0.95), PptInches(1.48), PptInches(11.4), PptInches(0.32), "Switch to the browser. Four clicks: home, Icehouse, publish a heading, then the CTA event ID.", size=13, bold=True, color=WHITE_PPT)
    demo = [
        (SHOTS / "home.png", "1", "Homepage", "Masthead is a section. Featured packet is a document — not a widget."),
        (SHOTS / "icehouse.png", "2", "Packet page", "Same Icehouse record. App Router, static params, ISR."),
        (SHOTS / "studio-home.png", "3", "Publish", "Studio → Catalog homepage → masthead → Publish → reload /."),
        (SHOTS / "order.png", "4", "Tracking", "“Request a packet list” is copy. eventId is the contract."),
    ]
    for i, (src, num, title, blurb) in enumerate(demo):
        col, row = i % 2, i // 2
        left = MARGIN + col * PptInches(6.1)
        top = PptInches(1.95) + row * PptInches(2.4)
        add_oval(s, left, top, PptInches(0.32), PptInches(0.32), MOSS_PPT)
        add_textbox(s, left, top + PptInches(0.04), PptInches(0.32), PptInches(0.26), num, size=11, bold=True, color=WHITE_PPT, align=PP_ALIGN.CENTER)
        add_textbox(s, left + PptInches(0.42), top, PptInches(5.4), PptInches(0.28), title, size=14, bold=True, color=INK_PPT, font=SERIF)
        add_browser_shot(s, src, left, top + PptInches(0.34), PptInches(5.85), PptInches(1.62), "")
        add_textbox(s, left, top + PptInches(2.0), PptInches(5.85), PptInches(0.32), blurb, size=11, color=MUTED_PPT)
    notes(s, BEATS[8])

    # 10 Decisions — list + mail-order screenshot
    s = content_slide(prs, "What I would defend", 10)
    decisions = [
        ("01", "RSC + GROQ", "Pages fetch on the server. Not a client-side CMS call on every view."),
        ("02", "References", "One variety document. Home, packet, and order form share it."),
        ("03", "eventId", "GTM keys off hero_order / nav_order. A rewrite of the button does not break measurement."),
        ("04", "ISR + webhook", "Revalidate about every 60s, or immediately after publish via a signed API."),
        ("05", "No cart", "A qualified packet-list request — a marketing lead, not a fake checkout."),
    ]
    for i, (num, title, body) in enumerate(decisions):
        top = PptInches(1.42) + i * PptInches(1.02)
        add_round(s, MARGIN, top, PptInches(5.55), PptInches(0.92), WHITE_PPT, RULE_PPT, 0.05)
        num_mark(s, PptInches(0.92), top + PptInches(0.24), num)
        add_textbox(s, PptInches(1.55), top + PptInches(0.1), PptInches(4.4), PptInches(0.32), title, size=15, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, PptInches(1.55), top + PptInches(0.44), PptInches(4.4), PptInches(0.4), body, size=11, color=MUTED_PPT)
    add_textbox(s, PptInches(6.5), PptInches(1.42), PptInches(6.1), PptInches(0.42), "We do not run a cart. Mark packets, tell us your zone, we write back — then pack the week it ships.", size=12, color=MUTED_PPT)
    add_browser_shot(s, SHOTS / "order.png", PptInches(6.5), PptInches(1.88), PptInches(6.1), PptInches(4.78), "/order  ·  eventId on the CTA")
    notes(s, BEATS[9])

    # 11 Close
    s = blank(prs, INK_PPT)
    add_rect(s, 0, 0, PptInches(0.1), SLIDE_H, MOSS_PPT)
    add_rect(s, 0, 0, SLIDE_W, PptInches(0.07), MOSS_PPT)
    add_textbox(s, MARGIN, PptInches(1.2), PptInches(11), PptInches(0.32), "CLOSE", size=12, bold=True, color=SAGE_PPT)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.85),
        PptInches(12),
        PptInches(2.9),
        "The catalog is the product.\nThe CMS is how it changes\nwithout a deploy.",
        size=36,
        bold=True,
        color=WHITE_PPT,
        font=SERIF,
        after=8,
    )
    add_rect(s, MARGIN, PptInches(5.15), PptInches(1.55), PptInches(0.07), MOSS_PPT)
    add_textbox(s, MARGIN, PptInches(5.4), PptInches(11), PptInches(0.4), "Two more live sites next. Public business only.", size=16, color=CREAM_PPT)
    add_textbox(s, MARGIN, PptInches(6.55), PptInches(8.5), PptInches(0.35), "Then questions", size=16, bold=True, color=WHITE_PPT, font=SERIF)
    add_textbox(s, PptInches(9.2), PptInches(6.55), PptInches(3.4), PptInches(0.35), "11  /  17", size=14, color=SAGE_PPT, align=PP_ALIGN.RIGHT)
    notes(s, BEATS[10])

    # 12 NDA map — public sites only
    s = content_slide(prs, "Work I cannot open", 12)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.38),
        PptInches(11.9),
        PptInches(0.42),
        "Same as an ads business: who buys, what is inventory, how a visit becomes a sale. Then where Sanity sits. No Studio, schemas, or fields.",
        size=14,
        color=MUTED_PPT,
    )
    nda = [
        {
            "left": MARGIN,
            "name": "Jamb  ·  jamb.co.uk",
            "shot": SHOTS / "jamb.png",
            "url": "jamb.co.uk  ·  public catalog",
            "lines": [
                "Pimlico dealer. Designers and collectors. Antiques plus reproductions.",
                "Next: who buys, inventory, the four collections.",
            ],
        },
        {
            "left": PptInches(6.78),
            "name": "Slingshot Bio  ·  slingshotbio.com",
            "shot": SHOTS / "slingshot.png",
            "url": "slingshotbio.com  ·  public catalog",
            "lines": [
                "Reagents for flow cytometry labs. Cell mimics, not donor cells.",
                "Next: the problem, the product, the shop.",
            ],
        },
    ]
    for card in nda:
        left = card["left"]
        add_round(s, left, PptInches(1.88), PptInches(5.85), PptInches(4.78), WHITE_PPT, RULE_PPT, 0.04)
        add_textbox(s, left + PptInches(0.22), PptInches(1.98), PptInches(5.4), PptInches(0.28), card["name"], size=14, bold=True, color=INK_PPT, font=SERIF)
        add_browser_shot(s, card["shot"], left + PptInches(0.22), PptInches(2.3), PptInches(5.4), PptInches(3.35), card["url"])
        add_lines(
            s,
            left + PptInches(0.22),
            PptInches(5.78),
            PptInches(5.4),
            PptInches(0.72),
            card["lines"],
            size=12,
            color=MUTED_PPT,
            after=4,
        )
    notes(s, BEATS[11])

    # 13 Jamb business
    s = content_slide(prs, "Jamb — who buys, what sells", 13)
    rows = [
        ("Who", "Interior designers, collectors, country-house clients. Flagship on Pimlico Road; showrooms in Los Angeles, Chicago, Dallas, Atlanta, Palm Beach."),
        ("Inventory", "Unique antiques — chimneypieces, lighting, furniture — and handmade reproductions from the south London workshop, so a design stays available after the antique sells."),
        ("Sale", "High-ticket, consultative. The site is the catalog for the showroom. A designer specs a piece; enquiry closes it. Not a cart."),
        ("CMS", "Sanity owns collections and journal stories. Next.js is the storefront. I cannot open the desk."),
    ]
    for i, (title, body) in enumerate(rows):
        top = PptInches(1.42) + i * PptInches(1.28)
        add_round(s, MARGIN, top, PptInches(5.45), PptInches(1.16), WHITE_PPT, RULE_PPT, 0.05)
        add_textbox(s, PptInches(0.95), top + PptInches(0.08), PptInches(4.9), PptInches(0.28), title, size=15, bold=True, color=MOSS_PPT, font=SERIF)
        add_textbox(s, PptInches(0.95), top + PptInches(0.38), PptInches(4.9), PptInches(0.7), body, size=12, color=MUTED_PPT)
    add_textbox(s, PptInches(6.38), PptInches(1.42), PptInches(6.2), PptInches(0.28), "Public homepage — not Studio", size=12, bold=True, color=MOSS_PPT)
    add_browser_shot(s, SHOTS / "jamb.png", PptInches(6.38), PptInches(1.72), PptInches(6.22), PptInches(4.94), "jamb.co.uk")
    notes(s, BEATS[12])

    # 14 Jamb catalog
    s = content_slide(prs, "Jamb — four public collections", 14)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.38),
        PptInches(11.9),
        PptInches(0.36),
        "Same contract as any marketing catalog: collections and stories in the CMS, storefront in Next.js. I will not name schemas or fields.",
        size=13,
        color=MUTED_PPT,
    )
    collections = [
        (SHOTS / "jamb-fireplaces.png", "01", "/fireplaces", "Fireplaces", "Antique chimneypieces and stone or marble reproductions. Bespoke when the room needs a scale that is not in stock."),
        (SHOTS / "jamb-lighting.png", "02", "/lighting", "Lighting", "Hanging globes, lanterns, wall lights. The Original Globe is the story that started the reproduction line."),
        (SHOTS / "jamb-furniture.png", "03", "/furniture", "Furniture", "English country-house seating and tables — antique and made. Same aesthetic as the chimneypieces."),
        (SHOTS / "jamb-journal.png", "04", "/journal", "Journal", "Stories for architects and designers. Not a blog calendar. Sanity owns the stories."),
    ]
    for i, (src, num, route, title, blurb) in enumerate(collections):
        col, row = i % 2, i // 2
        left = MARGIN + col * PptInches(6.1)
        top = PptInches(1.78) + row * PptInches(2.48)
        add_textbox(s, left, top, PptInches(0.7), PptInches(0.24), num, size=11, bold=True, color=MOSS_PPT)
        add_textbox(s, left + PptInches(0.7), top, PptInches(2.6), PptInches(0.24), title, size=13, bold=True, color=INK_PPT, font=SERIF)
        add_textbox(s, left + PptInches(3.4), top, PptInches(2.4), PptInches(0.24), route, size=11, color=MUTED_PPT, align=PP_ALIGN.RIGHT)
        add_browser_shot(s, src, left, top + PptInches(0.26), PptInches(5.85), PptInches(1.72), "")
        add_textbox(s, left, top + PptInches(2.02), PptInches(5.85), PptInches(0.38), blurb, size=11, color=MUTED_PPT)
    notes(s, BEATS[13])

    # 15 Slingshot business
    s = content_slide(prs, "Slingshot — the control problem", 15)
    rows = [
        ("Who", "Flow cytometry labs: biopharma, CROs, academic cores, cell therapy, instrument makers."),
        ("Problem", "Donor cells expire and drift lot to lot. Polystyrene beads have the wrong scatter. The control is the weak point in the assay."),
        ("Product", "Shelf-stable synthetic cell mimics — polymer particles that scatter, fluoresce, and stain like real cells. Catalog SKUs plus custom mimics."),
        ("Sale", "A scientist finds a control, reads a protocol, orders a reagent. Repeat purchases. Sanity owns product and resource documents."),
    ]
    for i, (title, body) in enumerate(rows):
        top = PptInches(1.42) + i * PptInches(1.28)
        add_round(s, MARGIN, top, PptInches(5.45), PptInches(1.16), WHITE_PPT, RULE_PPT, 0.05)
        add_textbox(s, PptInches(0.95), top + PptInches(0.08), PptInches(4.9), PptInches(0.28), title, size=15, bold=True, color=MOSS_PPT, font=SERIF)
        add_textbox(s, PptInches(0.95), top + PptInches(0.38), PptInches(4.9), PptInches(0.7), body, size=12, color=MUTED_PPT)
    add_textbox(s, PptInches(6.38), PptInches(1.42), PptInches(6.2), PptInches(0.28), "How Mimics Work — public argument", size=12, bold=True, color=MOSS_PPT)
    add_browser_shot(s, SHOTS / "slingshot-mimics.png", PptInches(6.38), PptInches(1.72), PptInches(6.22), PptInches(4.94), "slingshotbio.com/how-mimics-work")
    notes(s, BEATS[14])

    # 16 Slingshot catalog
    s = content_slide(prs, "Slingshot — catalog, argument, shop", 16)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.38),
        PptInches(11.9),
        PptInches(0.36),
        "Public surfaces only. Products and resources live in Sanity. Next.js is the catalog and shop. I will not open Studio.",
        size=13,
        color=MUTED_PPT,
    )
    sling_surfaces = [
        (SHOTS / "slingshot.png", "01", "/", "Homepage", "Catalog of controls: compensation, unmixing, viability, immunophenotyping."),
        (SHOTS / "slingshot-mimics.png", "02", "/how-mimics-work", "Argument", "Not a biologic. Not a bead. The page that sells the category."),
        (SHOTS / "slingshot-shop.png", "03", "/shop", "Shop", "SKUs by purpose, then a cart. Custom mimics on request. Repeat reagent purchases."),
    ]
    for i, (src, num, route, title, blurb) in enumerate(sling_surfaces):
        left = MARGIN + i * PptInches(4.05)
        add_textbox(s, left, PptInches(1.82), PptInches(0.5), PptInches(0.24), num, size=11, bold=True, color=MOSS_PPT)
        add_textbox(s, left + PptInches(0.5), PptInches(1.82), PptInches(3.3), PptInches(0.24), title, size=13, bold=True, color=INK_PPT, font=SERIF)
        add_browser_shot(s, src, left, PptInches(2.12), PptInches(3.9), PptInches(3.85), "")
        add_textbox(s, left, PptInches(6.05), PptInches(3.9), PptInches(0.72), f"{route}  ·  {blurb}", size=11, color=MUTED_PPT)
    notes(s, BEATS[15])

    # 17 Stop
    s = blank(prs, INK_PPT)
    add_rect(s, 0, 0, PptInches(0.1), SLIDE_H, MOSS_PPT)
    add_rect(s, 0, 0, SLIDE_W, PptInches(0.07), MOSS_PPT)
    add_textbox(s, MARGIN, PptInches(1.2), PptInches(11), PptInches(0.32), "STOP", size=12, bold=True, color=SAGE_PPT)
    add_textbox(
        s,
        MARGIN,
        PptInches(1.85),
        PptInches(12),
        PptInches(2.4),
        "I cannot open those desks.\nThis catalog I can.",
        size=36,
        bold=True,
        color=WHITE_PPT,
        font=SERIF,
        after=8,
    )
    add_rect(s, MARGIN, PptInches(5.15), PptInches(1.55), PptInches(0.07), MOSS_PPT)
    add_textbox(s, MARGIN, PptInches(5.4), PptInches(11), PptInches(0.4), "Schema, caching, and the order form — here, on Thorn & Furrow.", size=16, color=CREAM_PPT)
    add_textbox(s, MARGIN, PptInches(6.55), PptInches(8.5), PptInches(0.35), "Questions", size=16, bold=True, color=WHITE_PPT, font=SERIF)
    add_textbox(s, PptInches(9.2), PptInches(6.55), PptInches(3.4), PptInches(0.35), "17  /  17", size=14, color=SAGE_PPT, align=PP_ALIGN.RIGHT)
    notes(s, BEATS[16])

    path = Path(path) if path else OUT / "Thorn-and-Furrow-Portfolio.pptx"
    prs.save(path)
    return path


def build_script_md():
    lines = [
        "# Thorn & Furrow — 10-minute speaking script",
        "",
        "Read **SAY** out loud. Follow **DO**. Ten minutes is slides 1–11. Slides 12–17 if they have time.",
        "",
        "## Before you start",
        "",
        "- Catalog: `http://localhost:3000`",
        "- Studio: `http://localhost:3000/studio` (logged in)",
        "- Hard-refresh once. If Studio fails, keep going on the public site.",
        "- If you run long, skip slide 10 (decisions), close, and take questions. Slides 12–17 are extra if they ask about past work.",
        "",
        "## Clock",
        "",
        "| Time | Slide | You are doing |",
        "|---|---|---|",
        "| 0:00–0:25 | 1 Title | Marketing site, not SaaS. This catalog I can open. |",
        "| 0:25–1:05 | 2 Not a SaaS | Login app vs public catalog. |",
        "| 1:05–1:50 | 3 The stack | Next.js, Sanity, Vercel. |",
        "| 1:50–2:25 | 4 The job | Editors vs engineering. |",
        "| 2:25–3:00 | 5 Product | Four surfaces. No cart. |",
        "| 3:00–3:35 | 6 Why this shape | Documents, composition, event IDs. |",
        "| 3:35–4:25 | 7 Architecture | Studio → GROQ → Next.js → ISR. |",
        "| 4:25–5:00 | 8 Content model | Five documents. |",
        "| 5:00–7:40 | 9 Live demo | Home, packet, publish, tracking. |",
        "| 7:40–8:20 | 10 Decisions | Five calls you would make again. |",
        "| 8:20–8:50 | 11 Close | One line. Stop here if they only have ten minutes. |",
        "| 8:50–9:05 | 12 NDA map | Jamb and Slingshot. Public only. |",
        "| 9:05–9:40 | 13 Jamb business | Who, inventory, sale. |",
        "| 9:40–10:10 | 14 Jamb catalog | Fireplaces, lighting, furniture, journal. |",
        "| 10:10–10:45 | 15 Slingshot business | Problem, product, who pays. |",
        "| 10:45–11:15 | 16 Slingshot catalog | Home, how it works, shop. |",
        "| 11:15–11:30 | 17 Stop | I cannot open those desks. Questions here. |",
        "",
    ]
    for beat in BEATS:
        lines += [
            f"## {beat['clock']}  ·  {beat['slide']}",
            "",
        ]
        if beat.get("do"):
            lines += [f"**DO:** {beat['do']}", ""]
        lines += ["**SAY:**", "", beat["say"], ""]
    lines += [
        "## Close line (memorize)",
        "",
        "> The catalog is the product. The CMS is how it changes without a deploy.",
        "",
    ]
    path = OUT / "SCRIPT.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


if __name__ == "__main__":
    try:
        print(build_docx())
    except PermissionError:
        print("skipped docx (file is open)")
    try:
        print(build_pptx())
    except PermissionError:
        alt = OUT / "Thorn-and-Furrow-Deck.pptx"
        print(build_pptx(alt))
    print(build_script_md())
